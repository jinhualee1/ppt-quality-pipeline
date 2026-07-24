from __future__ import annotations

import json
import os
import re
import shutil
import subprocess
import sys
from dataclasses import asdict
from pathlib import Path
from typing import Any

from .io import relative_posix, write_json
from .models import RenderedDeck
from .text import extract_pptx_text


class RenderError(RuntimeError):
    pass


def _find_executable(env_name: str, candidates: list[str]) -> str | None:
    configured = os.environ.get(env_name, "").strip()
    if configured and Path(configured).is_file():
        return configured
    for candidate in candidates:
        resolved = shutil.which(candidate)
        if resolved:
            return resolved
        if Path(candidate).is_file():
            return candidate
    return None


def _node_executable() -> str | None:
    return _find_executable("PQP_NODE", ["node", r"C:\Program Files\nodejs\node.exe"])


def _powershell_executable() -> str | None:
    return _find_executable(
        "PQP_POWERSHELL",
        ["powershell.exe", "pwsh", r"C:\Windows\System32\WindowsPowerShell\v1.0\powershell.exe"],
    )


def libreoffice_executable() -> str | None:
    return _find_executable(
        "PQP_LIBREOFFICE",
        [
            "soffice",
            "libreoffice",
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
            "/usr/bin/libreoffice",
            "/usr/bin/soffice",
            "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        ],
    )


def pdftoppm_executable() -> str | None:
    return _find_executable("PQP_PDFTOPPM", ["pdftoppm", "pdftoppm.exe"])


def pdftoppm_available() -> bool:
    executable = pdftoppm_executable()
    if not executable:
        return False
    try:
        completed = subprocess.run(
            [executable, "-v"],
            capture_output=True,
            text=True,
            timeout=10,
            check=False,
        )
    except (OSError, subprocess.SubprocessError):
        return False
    return completed.returncode == 0


def powerpoint_available() -> bool:
    if sys.platform != "win32" or not _powershell_executable():
        return False
    try:
        import winreg

        with winreg.OpenKey(winreg.HKEY_CLASSES_ROOT, r"PowerPoint.Application\CLSID"):
            return True
    except (ImportError, FileNotFoundError, OSError):
        return False


def _project_root() -> Path:
    return Path(__file__).resolve().parents[2]


def _packaged_helper(folder: str, name: str) -> Path:
    packaged = Path(__file__).resolve().parent / folder / name
    source = _project_root() / "scripts" / name
    return packaged if packaged.is_file() else source


def render_html(source: Path, output_dir: Path, run_dir: Path) -> RenderedDeck:
    node = _node_executable()
    if not node:
        raise RenderError("Node.js was not found. Set PQP_NODE or install Node.js.")
    helper = _packaged_helper("browser", "render_html.mjs")
    if not helper.is_file():
        raise RenderError(f"HTML renderer helper is missing: {helper}")
    output_dir.mkdir(parents=True, exist_ok=True)
    command = [node, str(helper), "--input", str(source), "--output", str(output_dir)]
    completed = subprocess.run(
        command,
        capture_output=True,
        text=True,
        encoding="utf-8",
        timeout=180,
        check=False,
    )
    if completed.returncode != 0:
        detail = (completed.stderr or completed.stdout).strip()
        raise RenderError(detail or "HTML renderer failed.")
    manifest = json.loads((output_dir / "render_manifest.json").read_text(encoding="utf-8"))
    pages = [relative_posix(output_dir / name, run_dir) for name in manifest["pages"]]
    return RenderedDeck(
        artifact_path=relative_posix(source, run_dir),
        kind="html",
        pages=pages,
        page_count=len(pages),
        status="rendered",
        renderer="playwright",
        fidelity="high",
    )


def render_image(source: Path, output_dir: Path, run_dir: Path) -> RenderedDeck:
    output_dir.mkdir(parents=True, exist_ok=True)
    destination = output_dir / f"page_001{source.suffix.lower()}"
    shutil.copy2(source, destination)
    return RenderedDeck(
        artifact_path=relative_posix(source, run_dir),
        kind="image",
        pages=[relative_posix(destination, run_dir)],
        page_count=1,
        status="rendered",
        renderer="image-copy",
        fidelity="high",
    )


def _natural_number(path: Path) -> int:
    match = re.search(r"(\d+)(?!.*\d)", path.stem)
    return int(match.group(1)) if match else 0


def _render_pdf_pages(source: Path, output_dir: Path) -> tuple[list[Path], str]:
    output_dir.mkdir(parents=True, exist_ok=True)
    try:
        import fitz  # type: ignore[import-not-found]
    except ImportError:
        fitz = None

    if fitz is not None:
        pages = []
        with fitz.open(source) as document:
            for index, page in enumerate(document):
                pixmap = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
                destination = output_dir / f"page_{index + 1:03d}.png"
                pixmap.save(destination)
                pages.append(destination)
        return pages, "pymupdf"

    pdftoppm = pdftoppm_executable()
    if not pdftoppm:
        raise RenderError(
            "PDF rendering requires PyMuPDF or pdftoppm. Install the 'pdf' extra or set PQP_PDFTOPPM."
        )
    prefix = output_dir / "_page"
    completed = subprocess.run(
        [pdftoppm, "-png", "-r", "180", str(source), str(prefix)],
        capture_output=True,
        text=True,
        encoding="utf-8",
        errors="replace",
        timeout=300,
        check=False,
    )
    if completed.returncode != 0:
        detail = (completed.stderr or completed.stdout).strip()
        raise RenderError(detail or "pdftoppm failed.")
    generated = sorted(output_dir.glob("_page-*.png"), key=_natural_number)
    pages = []
    for index, generated_page in enumerate(generated, 1):
        destination = output_dir / f"page_{index:03d}.png"
        generated_page.replace(destination)
        pages.append(destination)
    if not pages:
        raise RenderError("PDF renderer produced no pages.")
    return pages, "pdftoppm"


def render_pdf(source: Path, output_dir: Path, run_dir: Path) -> RenderedDeck:
    pages, backend = _render_pdf_pages(source, output_dir)
    return RenderedDeck(
        artifact_path=relative_posix(source, run_dir),
        kind="pdf",
        pages=[relative_posix(page, run_dir) for page in pages],
        page_count=len(pages),
        status="rendered",
        renderer=backend,
        fidelity="high",
    )


def render_pptx_powerpoint(source: Path, output_dir: Path, run_dir: Path) -> RenderedDeck:
    if sys.platform != "win32":
        raise RenderError("Microsoft PowerPoint rendering is available only on Windows.")
    powershell = _powershell_executable()
    if not powershell:
        raise RenderError("PowerShell was not found.")
    helper = _packaged_helper("backends", "render_pptx_powerpoint.ps1")
    if not helper.is_file():
        raise RenderError(f"PowerPoint renderer helper is missing: {helper}")
    output_dir.mkdir(parents=True, exist_ok=True)
    command = [
        powershell,
        "-NoLogo",
        "-NoProfile",
        "-NonInteractive",
        "-ExecutionPolicy",
        "Bypass",
        "-File",
        str(helper),
        "-InputPath",
        str(source.resolve()),
        "-OutputDir",
        str(output_dir.resolve()),
        "-Height",
        os.environ.get("PQP_PPTX_HEIGHT", "1080"),
    ]
    completed = subprocess.run(
        command,
        capture_output=True,
        text=True,
        encoding="utf-8",
        errors="replace",
        timeout=300,
        check=False,
    )
    if completed.returncode != 0:
        detail = (completed.stderr or completed.stdout).strip()
        raise RenderError(detail or "Microsoft PowerPoint rendering failed.")
    pages = sorted(output_dir.glob("page_*.png"), key=_natural_number)
    if not pages:
        raise RenderError("Microsoft PowerPoint produced no rendered pages.")
    return RenderedDeck(
        artifact_path=relative_posix(source, run_dir),
        kind="pptx",
        pages=[relative_posix(page, run_dir) for page in pages],
        page_count=len(pages),
        status="rendered",
        renderer="microsoft-powerpoint",
        fidelity="high",
    )


def render_pptx_libreoffice(source: Path, output_dir: Path, run_dir: Path) -> RenderedDeck:
    soffice = libreoffice_executable()
    if not soffice:
        raise RenderError("LibreOffice was not found. Set PQP_LIBREOFFICE or install LibreOffice.")
    conversion_dir = output_dir / "_conversion"
    profile_dir = conversion_dir / "profile"
    conversion_dir.mkdir(parents=True, exist_ok=True)
    profile_uri = profile_dir.resolve().as_uri()
    command = [
        soffice,
        "--headless",
        f"-env:UserInstallation={profile_uri}",
        "--convert-to",
        "pdf",
        "--outdir",
        str(conversion_dir),
        str(source.resolve()),
    ]
    completed = subprocess.run(
        command,
        capture_output=True,
        text=True,
        encoding="utf-8",
        errors="replace",
        timeout=300,
        check=False,
    )
    if completed.returncode != 0:
        detail = (completed.stderr or completed.stdout).strip()
        raise RenderError(detail or "LibreOffice conversion failed.")
    pdfs = sorted(conversion_dir.glob("*.pdf"))
    if not pdfs:
        raise RenderError("LibreOffice produced no PDF.")
    try:
        pages, pdf_backend = _render_pdf_pages(pdfs[0], output_dir)
    finally:
        shutil.rmtree(conversion_dir, ignore_errors=True)
    return RenderedDeck(
        artifact_path=relative_posix(source, run_dir),
        kind="pptx",
        pages=[relative_posix(page, run_dir) for page in pages],
        page_count=len(pages),
        status="rendered",
        renderer=f"libreoffice+{pdf_backend}",
        fidelity="high",
    )


def _pptx_slide_texts(source: Path) -> list[str]:
    try:
        from pptx import Presentation  # type: ignore[import-not-found]
    except ImportError:
        text = extract_pptx_text(source)
        return [text] if text else []

    deck = Presentation(source)
    slide_texts = []
    for slide in deck.slides:
        chunks = [
            str(shape.text).strip()
            for shape in slide.shapes
            if hasattr(shape, "text") and str(shape.text).strip()
        ]
        slide_texts.append("\n".join(chunks))
    return slide_texts


def render_pptx_preview(
    source: Path,
    output_dir: Path,
    run_dir: Path,
    warnings: list[str] | None = None,
) -> RenderedDeck:
    try:
        from PIL import Image, ImageDraw, ImageFont
    except ImportError as exc:
        raise RenderError("PPTX fallback rendering requires Pillow.") from exc

    slide_texts = _pptx_slide_texts(source)
    if not slide_texts:
        raise RenderError("No slides could be read from the PPTX file.")
    output_dir.mkdir(parents=True, exist_ok=True)
    pages = []
    font = ImageFont.load_default()
    for index, text in enumerate(slide_texts):
        image = Image.new("RGB", (1280, 720), "#f7f8fa")
        draw = ImageDraw.Draw(image)
        draw.rectangle((0, 0, 1280, 70), fill="#17212b")
        draw.text((42, 25), f"Slide {index + 1}", fill="#ffffff", font=font)
        wrapped = []
        for paragraph in (text or "(No extractable text)").splitlines():
            while len(paragraph) > 92:
                wrapped.append(paragraph[:92])
                paragraph = paragraph[92:]
            wrapped.append(paragraph)
        draw.multiline_text((56, 110), "\n".join(wrapped[:22]), fill="#24313d", font=font, spacing=12)
        draw.text((56, 676), "Low-fidelity text preview", fill="#a15c05", font=font)
        destination = output_dir / f"page_{index + 1:03d}.png"
        image.save(destination)
        pages.append(relative_posix(destination, run_dir))
    return RenderedDeck(
        artifact_path=relative_posix(source, run_dir),
        kind="pptx",
        pages=pages,
        page_count=len(pages),
        status="rendered",
        renderer="pptx-text-preview",
        fidelity="low",
        warnings=warnings or ["No high-fidelity PPTX renderer was available."],
    )


def render_pptx(source: Path, output_dir: Path, run_dir: Path) -> RenderedDeck:
    backend = os.environ.get("PQP_PPTX_BACKEND", "auto").strip().lower()
    if backend not in {"auto", "powerpoint", "libreoffice", "preview"}:
        raise RenderError("PQP_PPTX_BACKEND must be auto, powerpoint, libreoffice, or preview.")

    errors = []
    if backend in {"auto", "powerpoint"}:
        try:
            return render_pptx_powerpoint(source, output_dir, run_dir)
        except RenderError as exc:
            errors.append(f"PowerPoint: {exc}")
            if backend == "powerpoint":
                raise

    if backend in {"auto", "libreoffice"}:
        try:
            return render_pptx_libreoffice(source, output_dir, run_dir)
        except RenderError as exc:
            errors.append(f"LibreOffice: {exc}")
            if backend == "libreoffice":
                raise

    if backend == "preview" or os.environ.get("PQP_PPTX_STRICT", "").lower() not in {"1", "true", "yes"}:
        return render_pptx_preview(source, output_dir, run_dir, errors)
    raise RenderError("; ".join(errors) or "No high-fidelity PPTX renderer was available.")


def render_artifact(artifact: dict[str, Any], run_dir: Path, output_dir: Path) -> RenderedDeck:
    source = run_dir / artifact["staged_path"]
    kind = artifact["kind"]
    try:
        if kind == "html":
            result = render_html(source, output_dir, run_dir)
        elif kind == "image":
            result = render_image(source, output_dir, run_dir)
        elif kind == "pdf":
            result = render_pdf(source, output_dir, run_dir)
        elif kind == "pptx":
            result = render_pptx(source, output_dir, run_dir)
        else:
            raise RenderError(f"No renderer is registered for artifact kind '{kind}'.")
    except (OSError, RenderError, subprocess.SubprocessError) as exc:
        result = RenderedDeck(
            artifact_path=relative_posix(source, run_dir),
            kind=kind,
            status="failed",
            renderer="",
            fidelity="none",
            error=str(exc),
        )
    write_json(output_dir / "deck.json", asdict(result))
    return result
